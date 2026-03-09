"""
CCRT Hematologic Toxicity - Dataset & ERD PPT Generator
"""
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = '/home/user/CCRT_Hematologic_Toxicity'

# ============================================================
# Data Loading
# ============================================================
df = pd.read_csv(f'{BASE}/data/raw/emr_synthetic/emr_patients.csv')
emr = pd.read_csv(f'{BASE}/data/processed/emr_processed.csv')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# Color Palette
# ============================================================
BG_DARK = RGBColor(0x1A, 0x23, 0x7E)      # Deep indigo
BG_MEDIUM = RGBColor(0x28, 0x35, 0x93)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
ACCENT1 = RGBColor(0x42, 0xA5, 0xF5)      # Blue
ACCENT2 = RGBColor(0xFF, 0x98, 0x00)      # Orange
ACCENT3 = RGBColor(0x4C, 0xAF, 0x50)      # Green
ACCENT4 = RGBColor(0xF4, 0x43, 0x36)      # Red
ACCENT5 = RGBColor(0x9C, 0x27, 0xB0)      # Purple
ACCENT6 = RGBColor(0x00, 0x96, 0x88)      # Teal
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xEC, 0xEF, 0xF1)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=11, color=BLACK, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name='Calibri'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Title decoration line
add_shape(slide, Inches(0.5), Inches(2.8), Inches(3), Pt(4), ACCENT1)

add_text_box(slide, Inches(0.5), Inches(3.0), Inches(12), Inches(1.5),
             'CCRT Hematologic Toxicity Prediction',
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.8),
             'Patient Dataset Overview & Project Architecture',
             font_size=24, color=ACCENT1, bold=False)

add_text_box(slide, Inches(0.5), Inches(5.5), Inches(8), Inches(0.5),
             'N=200 Synthetic EMR Cohort  |  LSTM + Temporal Attention  |  CTCAE v5.0 Grading',
             font_size=14, color=RGBColor(0xB0, 0xBE, 0xC5))

add_text_box(slide, Inches(0.5), Inches(6.5), Inches(4), Inches(0.4),
             'Concurrent Chemoradiation Therapy  |  Lung Cancer  |  Stage III',
             font_size=11, color=GRAY)

# ============================================================
# Slide 2: Dataset Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

# Header bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'Patient Dataset Summary (N=200)',
             font_size=28, color=WHITE, bold=True)

# Key metrics cards
metrics = [
    ('200', 'Total Patients', ACCENT1),
    (f'{df["age"].mean():.1f}', 'Mean Age (yrs)', ACCENT2),
    ('73% / 27%', 'Male / Female', ACCENT3),
    ('5', 'Chemo Regimens', ACCENT5),
    ('48', 'CBC Measures\n(6 vars x 8 weeks)', ACCENT6),
    (f'{emr["grade3_neutropenia"].mean()*100:.0f}%', 'Gr3+ Neutropenia', ACCENT4),
]

for i, (val, label, color) in enumerate(metrics):
    x = Inches(0.4 + i * 2.1)
    y = Inches(1.3)
    card = add_shape(slide, x, y, Inches(1.9), Inches(1.5), LIGHT_GRAY, color, 2)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(1.7), Inches(0.7),
                 val, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(1.7), Inches(0.5),
                 label, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# Data table
table_data = [
    ['Variable', 'Type', 'Values / Range', 'N features'],
    ['Demographics', 'Baseline', 'age, sex, bmi, ecog_ps', '4'],
    ['Cancer Stage', 'Baseline', 'stage (IIIA/B/C), t_stage (T1-4), n_stage (N0-3)', '3'],
    ['Lab Values', 'Baseline', f'creatinine ({df["creatinine"].min():.2f}-{df["creatinine"].max():.2f}), albumin ({df["albumin"].min():.1f}-{df["albumin"].max():.1f})', '2'],
    ['RT Treatment', 'Treatment', f'total_dose ({df["rt_total_dose"].min():.0f}-{df["rt_total_dose"].max():.0f}Gy), fraction_dose', '2'],
    ['Chemotherapy', 'Treatment', 'regimen (EP/TP/GP/DP/wPTX), dose, cycles (2-4)', '3'],
    ['CBC (raw)', 'Time-series', 'WBC, ANC, ALC, AMC, PLT, Hb  x  Week 0-7', '48'],
    ['CBC (derived)', 'Engineered', 'delta, slope, pct_change, cv, nadir, ratios', '~52'],
    ['Targets', 'Binary', 'grade3_neutropenia, anemia, thrombocytopenia, leukopenia, lymphopenia', '5'],
]

rows, cols = len(table_data), len(table_data[0])
tbl = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(3.1), Inches(12.5), Inches(3.8)).table

col_widths = [Inches(2.0), Inches(1.5), Inches(7.0), Inches(2.0)]
for j, w in enumerate(col_widths):
    tbl.columns[j].width = w

for i in range(rows):
    for j in range(cols):
        cell = tbl.cell(i, j)
        cell.text = table_data[i][j]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = 'Calibri'

        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_DARK
            p.font.color.rgb = WHITE
            p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
            p.font.color.rgb = BLACK

# ============================================================
# Slide 3: Demographics & Clinical Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'Patient Demographics & Clinical Characteristics',
             font_size=28, color=WHITE, bold=True)

# Age distribution text-based
age_box = add_shape(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(2.8), LIGHT_GRAY, ACCENT1, 1.5)
txBox = add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.4),
                     'Age Distribution', font_size=16, color=ACCENT1, bold=True)

tf = slide.shapes.add_textbox(Inches(0.6), Inches(1.85), Inches(5.8), Inches(2.0)).text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = f'Range: {df["age"].min()} - {df["age"].max()} years'
p.font.size = Pt(12); p.font.name = 'Calibri'

age_bins = pd.cut(df['age'], bins=[30,50,60,70,80,90])
for cat, cnt in age_bins.value_counts().sort_index().items():
    pct = cnt/len(df)*100
    bar = '█' * int(pct/2) + '░' * (25 - int(pct/2))
    add_paragraph(tf, f'{str(cat):>12s}  {bar}  {cnt:3d} ({pct:.0f}%)',
                  font_size=11, font_name='Consolas', color=BLACK)

add_paragraph(tf, f'\nMean: {df["age"].mean():.1f}  |  Median: {df["age"].median():.1f}  |  SD: {df["age"].std():.1f}',
              font_size=10, color=GRAY)

# Sex / ECOG
sex_box = add_shape(slide, Inches(6.9), Inches(1.3), Inches(3.0), Inches(2.8), LIGHT_GRAY, ACCENT3, 1.5)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(1.4), Inches(2.6), Inches(2.6))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Sex'; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT3
m = (df['sex']=='M').sum(); f_ = (df['sex']=='F').sum()
add_paragraph(tf, f'Male:    {m} ({m/len(df)*100:.0f}%)', font_size=13, bold=True, color=ACCENT1)
add_paragraph(tf, f'Female:  {f_} ({f_/len(df)*100:.0f}%)', font_size=13, bold=True, color=ACCENT4)
add_paragraph(tf, '', font_size=6)
add_paragraph(tf, 'ECOG PS', font_size=16, bold=True, color=ACCENT3, space_before=4)
for ps in sorted(df['ecog_ps'].unique()):
    cnt = (df['ecog_ps']==ps).sum()
    add_paragraph(tf, f'  PS {ps}: {cnt} ({cnt/len(df)*100:.0f}%)', font_size=12)

# BMI
bmi_box = add_shape(slide, Inches(10.2), Inches(1.3), Inches(2.8), Inches(2.8), LIGHT_GRAY, ACCENT2, 1.5)
txBox = slide.shapes.add_textbox(Inches(10.4), Inches(1.4), Inches(2.4), Inches(2.6))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'BMI'; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT2
add_paragraph(tf, f'Range: {df["bmi"].min():.1f}-{df["bmi"].max():.1f}', font_size=11)
add_paragraph(tf, f'Mean: {df["bmi"].mean():.1f}', font_size=11)
bmi_cats = pd.cut(df['bmi'], bins=[0,18.5,25,30,50], labels=['<18.5','18.5-25','25-30','>30'])
for cat in ['<18.5','18.5-25','25-30','>30']:
    cnt = (bmi_cats==cat).sum()
    add_paragraph(tf, f'  {cat}: {cnt} ({cnt/len(df)*100:.0f}%)', font_size=11)

# Stage / TNM
stage_box = add_shape(slide, Inches(0.4), Inches(4.4), Inches(4.0), Inches(2.8), LIGHT_GRAY, ACCENT5, 1.5)
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(3.6), Inches(2.6))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Cancer Stage (AJCC)'; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT5
for stg in ['IIIA','IIIB','IIIC']:
    cnt = (df['stage']==stg).sum()
    bar = '██' * int(cnt/len(df)*20)
    add_paragraph(tf, f'{stg}:  {bar}  {cnt} ({cnt/len(df)*100:.0f}%)',
                  font_size=12, font_name='Consolas')
add_paragraph(tf, '', font_size=4)
add_paragraph(tf, 'T-Stage:', font_size=11, bold=True)
t_str = '  '.join([f'T{t}: {(df["t_stage"]==f"T{t}").sum()}' for t in range(1,5)])
add_paragraph(tf, f'  {t_str}', font_size=10)
add_paragraph(tf, 'N-Stage:', font_size=11, bold=True)
n_str = '  '.join([f'N{n}: {(df["n_stage"]==f"N{n}").sum()}' for n in range(4)])
add_paragraph(tf, f'  {n_str}', font_size=10)

# Treatment
tx_box = add_shape(slide, Inches(4.7), Inches(4.4), Inches(4.2), Inches(2.8), LIGHT_GRAY, ACCENT6, 1.5)
txBox = slide.shapes.add_textbox(Inches(4.9), Inches(4.5), Inches(3.8), Inches(2.6))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Treatment'; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT6
add_paragraph(tf, 'RT Total Dose:', font_size=12, bold=True, space_before=2)
for dose in sorted(df['rt_total_dose'].unique()):
    cnt = (df['rt_total_dose']==dose).sum()
    add_paragraph(tf, f'  {dose:.0f} Gy: {cnt} ({cnt/len(df)*100:.0f}%)', font_size=11)
add_paragraph(tf, f'\nChemo Cycles: {df["chemo_cycles"].min()}-{df["chemo_cycles"].max()} (median: {df["chemo_cycles"].median():.0f})',
              font_size=11)

# Chemo Regimen
chemo_box = add_shape(slide, Inches(9.2), Inches(4.4), Inches(3.8), Inches(2.8), LIGHT_GRAY, ACCENT4, 1.5)
txBox = slide.shapes.add_textbox(Inches(9.4), Inches(4.5), Inches(3.4), Inches(2.6))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Chemotherapy Regimen'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT4
chemo_counts = df['chemo_regimen'].value_counts()
for reg, cnt in chemo_counts.items():
    pct = cnt/len(df)*100
    bar = '█' * int(pct/2)
    add_paragraph(tf, f'{reg:>18s}  {bar} {cnt} ({pct:.0f}%)',
                  font_size=11, font_name='Consolas')

# ============================================================
# Slide 4: CBC Time-series & Toxicity Outcomes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'CBC Time-Series Data & Hematologic Toxicity Outcomes',
             font_size=28, color=WHITE, bold=True)

# CBC Variables
cbc_box = add_shape(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(3.0), LIGHT_GRAY, ACCENT1, 1.5)
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.8))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'CBC Variables (6 measures x 8 timepoints)'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT1

cbc_vars = [
    ('WBC', '10^3/uL', 'White Blood Cell count'),
    ('ANC', '10^3/uL', 'Absolute Neutrophil Count'),
    ('ALC', '10^3/uL', 'Absolute Lymphocyte Count'),
    ('AMC', '10^3/uL', 'Absolute Monocyte Count'),
    ('PLT', '10^3/uL', 'Platelet count'),
    ('Hb',  'g/dL',    'Hemoglobin'),
]

for var, unit, desc in cbc_vars:
    w0_col = f'{var}_week0'
    vals = emr[w0_col]
    add_paragraph(tf, f'{var:>4s} ({unit:>7s}): {desc}  [{vals.min():.1f} - {vals.max():.1f}]',
                  font_size=11, font_name='Consolas')

add_paragraph(tf, '\nTimepoints: Week 0 (baseline) through Week 7', font_size=10, color=GRAY)
add_paragraph(tf, 'Input for prediction: Week 0, 1, 2 only', font_size=10, color=ACCENT4, bold=True)

# CBC Baseline Stats Table
cbc_stats_box = add_shape(slide, Inches(6.9), Inches(1.3), Inches(6.1), Inches(3.0), LIGHT_GRAY, ACCENT6, 1.5)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.7), Inches(0.35),
             'Baseline CBC Statistics (Week 0)', font_size=14, color=ACCENT6, bold=True)

stats_data = [['Variable', 'Mean', 'SD', 'Min', 'Median', 'Max']]
for var in ['WBC', 'ANC', 'ALC', 'AMC', 'PLT', 'Hb']:
    col = f'{var}_week0'
    v = emr[col]
    stats_data.append([var, f'{v.mean():.2f}', f'{v.std():.2f}', f'{v.min():.2f}',
                       f'{v.median():.2f}', f'{v.max():.2f}'])

tbl = slide.shapes.add_table(len(stats_data), 6, Inches(7.1), Inches(1.85), Inches(5.7), Inches(2.3)).table
for j, w in enumerate([Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9)]):
    tbl.columns[j].width = w

for i in range(len(stats_data)):
    for j in range(6):
        cell = tbl.cell(i, j)
        cell.text = stats_data[i][j]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10); p.font.name = 'Consolas'
        p.alignment = PP_ALIGN.CENTER
        if i == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT6
            p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Calibri'
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else LIGHT_GRAY

# Toxicity Outcomes
tox_box = add_shape(slide, Inches(0.4), Inches(4.6), Inches(6.2), Inches(2.6), LIGHT_GRAY, ACCENT4, 1.5)
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.4))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Grade 3+ Toxicity Outcomes (CTCAE v5.0)'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ACCENT4

tox_types = [
    ('grade3_neutropenia', 'Neutropenia', 'ANC < 1.0', ACCENT4),
    ('grade3_leukopenia', 'Leukopenia', 'WBC < 2.0', ACCENT1),
    ('grade3_lymphopenia', 'Lymphopenia', 'ALC < 0.5', ACCENT5),
    ('grade3_thrombocytopenia', 'Thrombocytopenia', 'PLT < 50', ACCENT2),
    ('grade3_anemia', 'Anemia', 'Hb < 8.0', ACCENT6),
]

for col, name, threshold, color in tox_types:
    rate = emr[col].mean() * 100
    n_pos = int(emr[col].sum())
    bar_len = int(rate / 2)
    bar = '█' * bar_len + '░' * (25 - bar_len)
    add_paragraph(tf, f'{name:>20s} ({threshold}):  {bar}  {rate:.1f}% (n={n_pos})',
                  font_size=10, font_name='Consolas')

# CTCAE Criteria
ctcae_box = add_shape(slide, Inches(6.9), Inches(4.6), Inches(6.1), Inches(2.6), LIGHT_GRAY, RGBColor(0x88,0x0E,0x4F), 1.5)
txBox = slide.shapes.add_textbox(Inches(7.1), Inches(4.7), Inches(5.7), Inches(2.4))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'CTCAE v5.0 Grading Criteria'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = RGBColor(0x88,0x0E,0x4F)

criteria = [
    'Neutropenia (ANC):   Grade 1: <2.0  |  Grade 2: <1.5  |  Grade 3: <1.0  |  Grade 4: <0.5',
    'Anemia (Hb):         Grade 1: <LLN  |  Grade 2: <10   |  Grade 3: <8.0  |  Grade 4: life-threatening',
    'Thrombocytopenia:    Grade 1: <LLN  |  Grade 2: <75   |  Grade 3: <50   |  Grade 4: <25',
    'Leukopenia (WBC):    Grade 1: <LLN  |  Grade 2: <3.0  |  Grade 3: <2.0  |  Grade 4: <1.0',
    'Lymphopenia (ALC):   Grade 1: <LLN  |  Grade 2: <0.8  |  Grade 3: <0.5  |  Grade 4: <0.2',
]
for c in criteria:
    add_paragraph(tf, c, font_size=9, font_name='Consolas')

add_paragraph(tf, '\nPrimary endpoint: Grade 3+ Neutropenia during CCRT (Week 0-7)',
              font_size=10, color=ACCENT4, bold=True)

# ============================================================
# Slide 5: Project Architecture ERD (Text-based)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'Project Architecture - ERD Overview',
             font_size=28, color=WHITE, bold=True)

# Insert the ERD image
erd_path = f'{BASE}/outputs/figures/project_erd.png'
if os.path.exists(erd_path):
    slide.shapes.add_picture(erd_path, Inches(0.3), Inches(1.15), Inches(12.7), Inches(6.1))

# ============================================================
# Slide 6: Patient Dataset Overview Image
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'Patient Dataset - Visual Dashboard',
             font_size=28, color=WHITE, bold=True)

overview_path = f'{BASE}/outputs/figures/patient_dataset_overview.png'
if os.path.exists(overview_path):
    slide.shapes.add_picture(overview_path, Inches(0.3), Inches(1.15), Inches(12.7), Inches(6.1))

# ============================================================
# Slide 7: Data Flow ERD (Detailed Text)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'Data Pipeline & Model Architecture',
             font_size=28, color=WHITE, bold=True)

# 4 layers
layers = [
    ('1. DATA LAYER', RGBColor(0xE6,0x51,0x00), [
        'emr_patients.csv (200 rows x 16 cols)',
        '  PK: patient_id | demographics + stage + treatment',
        'emr_cbc_results.csv (~600 rows, long format)',
        '  FK: patient_id | exam_date + 6 CBC variables',
        'ctcae_v5_criteria.json',
        '  Grade thresholds for 5 toxicity types',
    ]),
    ('2. PREPROCESSING LAYER', BG_MEDIUM, [
        'preprocessing.py: Date -> Week mapping, Long -> Wide pivot, CTCAE Grade calc',
        'data_loader.py: CSV loading, missing value imputation (median), outlier clipping (IQR)',
        'feature_engineer.py: Temporal features (delta, slope, cv, nadir, ratios)',
        'dataset.py: PyTorch CCRTDataset with 80% masking augmentation (Week 3-7)',
        'Output: emr_processed.csv (200 rows x 108 columns)',
    ]),
    ('3. MODEL LAYER', RGBColor(0x00,0x4D,0x40), [
        'Baseline Models: XGBoost + LightGBM + Logistic Regression',
        'LSTM: LSTMPredictor (2-layer LSTM + Temporal Attention + FC baseline encoder)',
        'Forecaster: CBCForecasterV2 (Seq2Seq: Week 0-2 -> predict Week 3-7)',
        'Training: 5-fold Stratified CV, EarlyStopping, SMOTE, class weights',
        'Entry: main.py (4 modes) | lstm_forecast.py | lstm_e2e.py | external_validation.py',
    ]),
    ('4. EVALUATION LAYER', RGBColor(0xF5,0x7F,0x17), [
        'metrics.py: AUROC, AUPRC, Sensitivity, Specificity, PPV, NPV, F1, Bootstrap CI',
        'visualization.py: ROC/PR curves, confusion matrix, feature importance, CBC plots',
        'helpers.py: Logging, seed setting, result serialization',
        'Output: 14+ figures, config.yaml, experiment_results.json',
    ]),
]

y_start = Inches(1.2)
for i, (title, color, items) in enumerate(layers):
    y = y_start + Inches(i * 1.55)
    # Title bar
    add_shape(slide, Inches(0.4), y, Inches(12.5), Inches(0.35), color)
    add_text_box(slide, Inches(0.6), y + Pt(1), Inches(12), Inches(0.35),
                 title, font_size=13, color=WHITE, bold=True)
    # Items
    txBox = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.4), Inches(12), Inches(1.1))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = items[0]; p.font.size = Pt(10); p.font.name = 'Consolas'
    for item in items[1:]:
        add_paragraph(tf, item, font_size=10, font_name='Consolas')

# Flow arrows (text)
add_text_box(slide, Inches(6.2), Inches(2.63), Inches(0.8), Inches(0.3),
             '▼', font_size=18, color=ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.2), Inches(4.18), Inches(0.8), Inches(0.3),
             '▼', font_size=18, color=ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.2), Inches(5.73), Inches(0.8), Inches(0.3),
             '▼', font_size=18, color=ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 8: Entity Relationship Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             'Entity Relationships & Key Schema',
             font_size=28, color=WHITE, bold=True)

# Entity boxes
entities = [
    ('emr_patients', Inches(0.4), Inches(1.3), Inches(3.8), Inches(4.5),
     RGBColor(0xE6,0x51,0x00), [
         'PK  patient_id  VARCHAR',
         '--- Demographics ---',
         '    age          INT',
         '    sex          M/F',
         '    bmi          FLOAT',
         '    ecog_ps      INT (0-4)',
         '--- Staging ---',
         '    stage        IIIA/B/C',
         '    t_stage      T1-T4',
         '    n_stage      N0-N3',
         '--- Lab ---',
         '    creatinine   FLOAT',
         '    albumin      FLOAT',
         '--- Treatment ---',
         '    rt_start_date    DATE',
         '    rt_total_dose    FLOAT',
         '    rt_fraction_dose FLOAT',
         '    chemo_regimen    VARCHAR',
         '    chemo_dose       FLOAT',
         '    chemo_cycles     INT',
     ]),
    ('emr_cbc_results', Inches(4.7), Inches(1.3), Inches(3.3), Inches(3.5),
     ACCENT1, [
         'FK  patient_id  VARCHAR',
         '    exam_date    DATE',
         '--- CBC Values ---',
         '    WBC          FLOAT',
         '    ANC          FLOAT',
         '    ALC          FLOAT',
         '    AMC          FLOAT',
         '    PLT          FLOAT',
         '    Hb           FLOAT',
         '',
         '~600 rows (long fmt)',
     ]),
    ('emr_processed', Inches(8.5), Inches(1.3), Inches(4.5), Inches(4.5),
     ACCENT3, [
         'PK  patient_id       VARCHAR',
         '--- CBC Wide Format ---',
         '    {var}_week{0..7}  FLOAT (48)',
         '--- CTCAE Grades ---',
         '    {tox}_grade_week{0..7} (40)',
         '    max_grade_{tox}       (5)',
         '--- Targets ---',
         '    grade3_neutropenia    BINARY',
         '    grade3_anemia         BINARY',
         '    grade3_thrombocytopenia BINARY',
         '    grade3_leukopenia     BINARY',
         '    grade3_lymphopenia    BINARY',
         '--- Clinical (merged) ---',
         '    age, sex, bmi, etc.   (16)',
         '',
         '200 rows x 108 columns',
     ]),
]

for name, x, y, w, h, color, fields in entities:
    # Header
    add_shape(slide, x, y, w, Inches(0.45), color)
    add_text_box(slide, x + Inches(0.15), y + Pt(2), w - Inches(0.3), Inches(0.4),
                 name, font_size=13, color=WHITE, bold=True)
    # Body
    body = add_shape(slide, x, y + Inches(0.45), w, h - Inches(0.45), LIGHT_GRAY, color, 1)
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), h - Inches(0.65))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = fields[0]; p.font.size = Pt(9); p.font.name = 'Consolas'
    if fields[0].startswith('PK') or fields[0].startswith('FK'):
        p.font.color.rgb = ACCENT4; p.font.bold = True
    for fld in fields[1:]:
        pp = add_paragraph(tf, fld, font_size=9, font_name='Consolas')
        if fld.startswith('PK') or fld.startswith('FK'):
            pp.font.color.rgb = ACCENT4; pp.font.bold = True
        elif fld.startswith('---'):
            pp.font.color.rgb = GRAY; pp.font.bold = True

# Relationship labels
add_text_box(slide, Inches(3.9), Inches(2.5), Inches(1.0), Inches(0.4),
             '1 : N  ──►', font_size=12, color=ACCENT4, bold=True)
add_text_box(slide, Inches(7.7), Inches(2.5), Inches(1.2), Inches(0.4),
             'merge ──►', font_size=12, color=ACCENT3, bold=True)

# Relationship description
rel_box = add_shape(slide, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.2),
                    RGBColor(0xFD,0xF2,0xE9), RGBColor(0xE6,0x51,0x00), 1)
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(12), Inches(1.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = 'Entity Relationships:'; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = BG_DARK
add_paragraph(tf, 'emr_patients (1) ──── (N) emr_cbc_results    via patient_id    [1 patient has multiple CBC exam records over 8 weeks]', font_size=10, font_name='Consolas')
add_paragraph(tf, 'emr_patients + emr_cbc_results  ──merge──►  emr_processed     [Long-to-Wide pivot + CTCAE grade calculation + feature merge]', font_size=10, font_name='Consolas')
add_paragraph(tf, 'emr_processed  ──split──►  train/val/test sets (80/10/10)       [Stratified by grade3_neutropenia to preserve class balance]', font_size=10, font_name='Consolas')

# ============================================================
# Save
# ============================================================
output_path = f'{BASE}/outputs/CCRT_Dataset_and_ERD.pptx'
prs.save(output_path)
print(f"\nPPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
